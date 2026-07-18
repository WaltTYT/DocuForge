"""核心转换引擎 - 基于 Pandoc + Python 库实现文档互转"""

import os
import subprocess
import tempfile
from pathlib import Path

# 支持的格式
SUPPORTED_FORMATS = {".docx", ".pdf", ".pptx", ".md"}

# Pandoc 直接支持的转换映射 (from -> to)
PANDOC_DIRECT = {
    (".md", ".docx"),
    (".md", ".pptx"),
    (".md", ".html"),
    (".docx", ".md"),
    (".docx", ".html"),
    (".html", ".md"),
    (".html", ".docx"),
}


def _find_pandoc() -> str:
    """查找 pandoc 可执行文件"""
    # 优先检查常见安装路径
    search_paths = [
        r"D:\study\pandoc\pandoc.exe",
        r"C:\Program Files\Pandoc\pandoc.exe",
        r"C:\Program Files (x86)\Pandoc\pandoc.exe",
    ]
    for p in search_paths:
        if os.path.isfile(p):
            return p
    # 回退到 PATH
    return "pandoc"


def _run_pandoc(input_file: str, output_file: str, extra_args: list = None):
    """调用 pandoc 执行转换"""
    pandoc = _find_pandoc()
    cmd = [pandoc, input_file, "-o", output_file]
    if extra_args:
        cmd.extend(extra_args)
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        raise RuntimeError(f"Pandoc 转换失败: {result.stderr.strip()}")


def _get_format(filepath: str) -> str:
    """获取文件扩展名（小写）"""
    return Path(filepath).suffix.lower()


def _check_format(fmt: str):
    """检查格式是否支持"""
    if fmt not in SUPPORTED_FORMATS:
        raise ValueError(f"不支持的格式: {fmt}，支持的格式: {', '.join(sorted(SUPPORTED_FORMATS))}")


# ─── PDF 读取 ───

def _extract_pdf_text(filepath: str) -> str:
    """使用 PyMuPDF 提取 PDF 文本"""
    import fitz
    doc = fitz.open(filepath)
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return "\n\n".join(pages)


def _extract_pdf_images(filepath: str, output_dir: str) -> list:
    """提取 PDF 中的图片"""
    import fitz
    doc = fitz.open(filepath)
    images = []
    for page_num in range(len(doc)):
        for img_index, img in enumerate(doc[page_num].get_images()):
            xref = img[0]
            pix = fitz.Pixmap(doc, xref)
            if pix.alpha:
                pix = fitz.Pixmap(doc, xref, alpha=False)
            img_path = os.path.join(output_dir, f"page{page_num+1}_img{img_index+1}.png")
            pix.save(img_path)
            images.append(img_path)
    doc.close()
    return images


# ─── PPTX 处理 ───

def _extract_pptx_content(filepath: str) -> str:
    """从 PPTX 提取文本内容为 Markdown"""
    from pptx import Presentation
    prs = Presentation(filepath)
    md_lines = []
    for i, slide in enumerate(prs.slides, 1):
        md_lines.append(f"## 幻灯片 {i}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        md_lines.append(text)
            if shape.has_table:
                table = shape.table
                md_lines.append("")
                # 表头
                header = [cell.text.strip() for cell in table.rows[0].cells]
                md_lines.append("| " + " | ".join(header) + " |")
                md_lines.append("| " + " | ".join(["---"] * len(header)) + " |")
                for row in list(table.rows)[1:]:
                    cells = [cell.text.strip() for cell in row.cells]
                    md_lines.append("| " + " | ".join(cells) + " |")
                md_lines.append("")
        md_lines.append("")
    return "\n".join(md_lines)


def _create_pptx_from_md(md_content: str, output_file: str):
    """从 Markdown 内容创建 PPTX（通过 pandoc）"""
    with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False, encoding="utf-8") as f:
        f.write(md_content)
        tmp_md = f.name
    try:
        _run_pandoc(tmp_md, output_file)
    finally:
        os.unlink(tmp_md)


def _create_pptx_from_docx(docx_path: str, output_file: str):
    """从 DOCX 转 PPTX：先转 md，再转 pptx"""
    with tempfile.NamedTemporaryFile(suffix=".md", delete=False) as f:
        tmp_md = f.name
    try:
        _run_pandoc(docx_path, tmp_md)
        _run_pandoc(tmp_md, output_file)
    finally:
        if os.path.isfile(tmp_md):
            os.unlink(tmp_md)


# ─── PDF 生成 ───

def _create_pdf_from_text(text: str, output_file: str):
    """从文本内容生成 PDF"""
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas as pdf_canvas
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    # 尝试注册中文字体
    font_name = "Helvetica"
    font_paths = [
        (r"C:\Windows\Fonts\msyh.ttc", "MSYH"),
        (r"C:\Windows\Fonts\simsun.ttc", "SimSun"),
        (r"C:\Windows\Fonts\simhei.ttf", "SimHei"),
    ]
    for fpath, fname in font_paths:
        if os.path.isfile(fpath):
            try:
                pdfmetrics.registerFont(TTFont(fname, fpath))
                font_name = fname
                break
            except Exception:
                continue

    c = pdf_canvas.Canvas(output_file, pagesize=A4)
    width, height = A4
    margin = 50
    y = height - margin
    line_height = 16

    for line in text.split("\n"):
        if y < margin:
            c.showPage()
            y = height - margin
        c.setFont(font_name, 11)
        c.drawString(margin, y, line)
        y -= line_height

    c.save()


def _create_pdf_from_md(md_path: str, output_file: str):
    """从 Markdown 生成 PDF（优先用 pandoc，回退到 reportlab）"""
    try:
        _run_pandoc(md_path, output_file, ["--pdf-engine=xelatex", "-V", "CJKmainfont=Microsoft YaHei"])
    except RuntimeError:
        # pandoc PDF 引擎不可用时，提取文本后用 reportlab 生成
        with open(md_path, "r", encoding="utf-8") as f:
            text = f.read()
        _create_pdf_from_text(text, output_file)


# ─── 公开 API ───

def convert(input_file: str, output_file: str):
    """
    将 input_file 转换为 output_file，格式由扩展名自动推断。

    支持的格式: .docx, .pdf, .pptx, .md
    """
    input_file = os.path.abspath(input_file)
    output_file = os.path.abspath(output_file)

    if not os.path.isfile(input_file):
        raise FileNotFoundError(f"输入文件不存在: {input_file}")

    src_fmt = _get_format(input_file)
    dst_fmt = _get_format(output_file)
    _check_format(src_fmt)
    _check_format(dst_fmt)

    if src_fmt == dst_fmt:
        raise ValueError("输入和输出格式相同，无需转换")

    os.makedirs(os.path.dirname(output_file) or ".", exist_ok=True)

    # 1. Pandoc 直接转换
    if (src_fmt, dst_fmt) in PANDOC_DIRECT:
        _run_pandoc(input_file, output_file)
        return

    # 2. 涉及 PDF 输入: 提取文本 -> 转为目标格式
    if src_fmt == ".pdf":
        text = _extract_pdf_text(input_file)
        if dst_fmt == ".md":
            with open(output_file, "w", encoding="utf-8") as f:
                f.write(text)
        elif dst_fmt == ".docx":
            with tempfile.NamedTemporaryFile(suffix=".md", delete=False, mode="w", encoding="utf-8") as f:
                f.write(text)
                tmp_md = f.name
            try:
                _run_pandoc(tmp_md, output_file)
            finally:
                os.unlink(tmp_md)
        elif dst_fmt == ".pptx":
            _create_pptx_from_md(text, output_file)
        return

    # 3. 涉及 PDF 输出
    if dst_fmt == ".pdf":
        if src_fmt == ".md":
            _create_pdf_from_md(input_file, output_file)
        elif src_fmt == ".docx":
            # docx -> md -> pdf
            with tempfile.NamedTemporaryFile(suffix=".md", delete=False) as f:
                tmp_md = f.name
            try:
                _run_pandoc(input_file, tmp_md)
                _create_pdf_from_md(tmp_md, output_file)
            finally:
                if os.path.isfile(tmp_md):
                    os.unlink(tmp_md)
        elif src_fmt == ".pptx":
            text = _extract_pptx_content(input_file)
            with tempfile.NamedTemporaryFile(suffix=".md", delete=False, mode="w", encoding="utf-8") as f:
                f.write(text)
                tmp_md = f.name
            try:
                _create_pdf_from_md(tmp_md, output_file)
            finally:
                os.unlink(tmp_md)
        return

    # 4. PPTX 与其他格式互转
    if src_fmt == ".pptx" and dst_fmt == ".md":
        text = _extract_pptx_content(input_file)
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(text)
        return

    if src_fmt == ".pptx" and dst_fmt == ".docx":
        text = _extract_pptx_content(input_file)
        _create_pptx_from_md(text, output_file) if False else None
        # 用 pandoc: md -> docx
        with tempfile.NamedTemporaryFile(suffix=".md", delete=False, mode="w", encoding="utf-8") as f:
            f.write(text)
            tmp_md = f.name
        try:
            _run_pandoc(tmp_md, output_file)
        finally:
            os.unlink(tmp_md)
        return

    if src_fmt == ".docx" and dst_fmt == ".pptx":
        _create_pptx_from_docx(input_file, output_file)
        return

    raise RuntimeError(f"不支持的转换路径: {src_fmt} -> {dst_fmt}")


def get_supported_formats() -> list:
    """返回支持的格式列表"""
    return sorted(SUPPORTED_FORMATS)
